from __future__ import annotations

import hashlib
from pathlib import Path


def _extract_pptx_images(prs, output_dir: Path) -> list[str]:
    """Extract embedded images from a PPTX presentation.

    Iterates over every slide's shapes and saves any picture shapes
    (shape.image.blob) to *output_dir*.  Returns a list of saved file paths.
    """
    image_paths: list[str] = []
    seen_hashes: set[str] = set()
    try:
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.shape_type or not hasattr(shape, "image"):
                    continue
                try:
                    blob = shape.image.blob
                    ct = shape.image.content_type or "image/png"
                except Exception:
                    continue
                if not blob:
                    continue
                ext = ct.split("/")[-1].replace("jpeg", "jpg")
                content_hash = hashlib.sha256(blob).hexdigest()[:12]
                if content_hash in seen_hashes:
                    continue
                seen_hashes.add(content_hash)
                filename = f"img_{content_hash}.{ext}"
                dest = output_dir / filename
                if not dest.exists():
                    dest.write_bytes(blob)
                image_paths.append(str(dest))
    except Exception:
        pass  # graceful — image extraction is best-effort
    return image_paths


def convert_pptx(path: str, images_dir: str | None = None) -> tuple[str, list[str]]:
    try:
        from pptx import Presentation
    except Exception as exc:
        raise RuntimeError("python-pptx is required for PPTX conversion") from exc

    file_path = Path(path)
    prs = Presentation(str(file_path))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)

    # Extract embedded images when an output directory is provided
    image_paths: list[str] = []
    if images_dir:
        out = Path(images_dir)
        out.mkdir(parents=True, exist_ok=True)
        image_paths = _extract_pptx_images(prs, out)

    return "\n".join(parts), image_paths
